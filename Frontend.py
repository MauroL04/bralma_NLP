import os
os.environ["TOKENIZERS_PARALLELISM"] = "false"

import streamlit as st
from pathlib import Path
import PyPDF2
from pptx import Presentation
from datetime import datetime
import warnings

# Import direct LLM modules + RAG
import groq_answer_llm
import langchain_agent

warnings.filterwarnings("ignore", category=SyntaxWarning, module="pysbd")
warnings.filterwarnings("ignore")
 
 
# Page configuration
st.set_page_config(
    page_title="Document Chat Assistant",
    page_icon="💬",
    layout="wide",
    initial_sidebar_state="expanded"
)
 
# Custom CSS for chatbot design
st.markdown("""
    <style>
    /* Sidebar styling */
    [data-testid="stSidebar"] {
        background-color: #1e1e1e;
        border-right: 1px solid #404040;
    }
   
    [data-testid="stSidebar"] [data-testid="stMarkdownContainer"] {
        color: #e0e0e0;
    }
   
    /* Main container styling */
    .main {
        background-color: #2d2d2d;
        color: white;
    }
   
    .block-container {
        padding-top: 3rem;
        padding-bottom: 3rem;
        max-width: 900px;
    }
   
    /* Hide default streamlit elements */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
   
    /* Chat message styling */
    .chat-message {
        padding: 1.2rem;
        border-radius: 0.8rem;
        margin-bottom: 1rem;
        animation: fadeIn 0.3s ease-in;
    }
   
    .chat-message.user {
        background-color: #3a3a3a;
        margin-left: 20%;
    }
   
    .chat-message.bot {
        background-color: #1e1e1e;
        margin-right: 20%;
        border: 1px solid #404040;
    }
   
    .message-content {
        color: #e0e0e0;
        line-height: 1.6;
    }
   
    .message-header {
        font-weight: 600;
        margin-bottom: 0.5rem;
        color: #a0a0a0;
        font-size: 0.85rem;
    }
   
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(10px); }
        to { opacity: 1; transform: translateY(0); }
    }
   
    /* Chat input styling */
    .stChatInput {
        position: fixed;
        bottom: 2rem;
        left: 50%;
        transform: translateX(-50%);
        width: 70%;
        max-width: 800px;
        background-color: transparent;
        padding: 0;
        z-index: 100;
    }
   
    .stChatInput > div {
        background-color: #f5f5f5;
        border-radius: 2rem;
        border: 2px solid #e0e0e0;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.3);
        transition: all 0.3s ease;
    }
   
    .stChatInput > div:hover {
        border-color: #4a90e2;
        box-shadow: 0 6px 16px rgba(74, 144, 226, 0.3);
    }
   
    .stChatInput input {
        color: #333 !important;
        font-size: 1rem;
    }
   
    .stChatInput input::placeholder {
        color: #999 !important;
    }
   
    /* File upload overlay */
    .upload-overlay {
        position: fixed;
        top: 0;
        left: 0;
        right: 0;
        bottom: 0;
        background-color: rgba(0, 0, 0, 0.9);
        display: none;
        justify-content: center;
        align-items: center;
        z-index: 9999;
        backdrop-filter: blur(5px);
    }
   
    .upload-box {
        border: 3px dashed #666;
        border-radius: 1rem;
        padding: 3rem;
        text-align: center;
        color: white;
        font-size: 1.5rem;
    }
   
    /* File status badge */
    .file-badge {
        position: fixed;
        top: 1rem;
        right: 1rem;
        background-color: #3a3a3a;
        padding: 0.5rem 1rem;
        border-radius: 2rem;
        font-size: 0.9rem;
        color: #4CAF50;
        border: 1px solid #4CAF50;
        z-index: 1000;
    }
   
    /* Title styling */
    h1 {
        color: white !important;
        text-align: center;
        margin-bottom: 0.5rem;
    }
   
    /* Subtitle styling */
    .subtitle {
        text-align: center;
        color: #888;
        margin-bottom: 2rem;
        font-size: 1rem;
    }
   
    /* Streamlit elements color fix */
    .stMarkdown, p {
        color: #e0e0e0;
    }
   
    /* Info box styling */
    .stAlert {
        background-color: #1e1e1e;
        color: #e0e0e0;
        border: 1px solid #404040;
    }
    </style>
    """, unsafe_allow_html=True)
 
# Initialize session state
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'uploaded_documents' not in st.session_state:
    st.session_state.uploaded_documents = []  # List of dicts: [{"filename": "...", "text": "...", "type": "pdf/pptx"}]
 
def extract_text_from_pdf(pdf_file):
    """Extract text from uploaded PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        return f"Error extracting text: {str(e)}"

def extract_text_from_ppt(ppt_file):
    """Extract text from uploaded PowerPoint file"""
    try:
        prs = Presentation(ppt_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting text: {str(e)}"
 
def get_combined_documents_context():
    """Combine all uploaded documents (PDF and PPTX) into one context string"""
    if not st.session_state.uploaded_documents:
        return ""
   
    combined = ""
    for idx, doc in enumerate(st.session_state.uploaded_documents):
        doc_type = doc.get('type', 'document').upper()
        combined += f"\n\n=== {doc_type} {idx+1}: {doc['filename']} ===\n"
        combined += doc['text'][:3000]  # Limit each document to avoid token overflow
   
    return combined
 
def get_bot_response(user_question):
    """
    Team's aanpak: Direct answer_and_maybe_quiz call.
    Enhanced with RAG semantic search for BETTER answers.
    
    Workflow:
        1. Get document context (PDF + PPTX uploaded files)
        2. Get RAG context (semantic search from ChromaDB)
        3. Combine both contexts
        4. Call answer_and_maybe_quiz for final answer + quiz
    """
    # Step 1: Get uploaded documents context
    document_context = get_combined_documents_context()
    
    # Step 2: Get RAG context (semantic search - better relevance)
    try:
        rag_context = langchain_agent.get_context_for_question(user_question, k=4)
    except:
        rag_context = ""  # Fallback if no documents ingested in ChromaDB
    
    # Step 3: Combine contexts (RAG is usually better, but include uploaded docs too)
    combined_context = rag_context if rag_context else document_context
    
    # Step 4: Direct LLM call (team's approach)
    return groq_answer_llm.answer_and_maybe_quiz(user_question, combined_context)
 
# Sidebar for uploaded files overview
with st.sidebar:
    st.header("📚 Uploaded Documents")
   
    if st.session_state.uploaded_documents:
        st.markdown(f"**{len(st.session_state.uploaded_documents)} file(s) uploaded**")
        st.markdown("---")
       
        for idx, doc in enumerate(st.session_state.uploaded_documents):
            col1, col2 = st.columns([4, 1])
            with col1:
                icon = "📄" if doc.get('type') == 'pdf' else "📊"
                st.markdown(f"{icon} **{doc['filename']}**")
                st.caption(f"{len(doc['text'])} characters")
            with col2:
                if st.button("🗑️", key=f"remove_{idx}", help="Remove file"):
                    st.session_state.uploaded_documents.pop(idx)
                    st.rerun()
       
        st.markdown("---")
        if st.button("🗑️ Clear All", use_container_width=True):
            st.session_state.uploaded_documents = []
            st.session_state.messages = []
            st.rerun()
    else:
        st.info("No files uploaded yet")
    
    # Project Analysis Section (Commented out - not yet implemented)
    # st.markdown("---")
    # st.header("🔬 Project Analysis")
    # if st.button("📊 Analyze Codebase", use_container_width=True):
    #     st.info("Code analysis feature coming soon!")
 
# App header
st.title("What's on the agenda today?")
 
# File upload section - PDF and PPTX support
uploaded_file = st.file_uploader(
    "Drop your PDF or PPTX file here",
    type=['pdf', 'pptx'],
    help="Drag and drop a PDF or PowerPoint file to upload",
    label_visibility="collapsed"
)
 
# Process uploaded file (PDF or PPTX)
if uploaded_file is not None:
    # Check if file already exists
    existing_names = [doc['filename'] for doc in st.session_state.uploaded_documents]
   
    if uploaded_file.name not in existing_names:
        # Detect file type and extract text
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        with st.spinner(f"Processing {file_extension.upper()}..."):
            if file_extension == 'pdf':
                text = extract_text_from_pdf(uploaded_file)
                doc_type = 'pdf'
            elif file_extension == 'pptx':
                text = extract_text_from_ppt(uploaded_file)
                doc_type = 'pptx'
            else:
                st.error("Unsupported file type")
                text = ""
                doc_type = 'unknown'
            
            if text and not text.startswith("Error"):
                # Add to session state
                st.session_state.uploaded_documents.append({
                    "filename": uploaded_file.name,
                    "text": text,
                    "type": doc_type
                })
                
                # Ingest into ChromaDB for RAG (better semantic search)
                try:
                    langchain_agent.ingest_pdf(text.encode('utf-8'), uploaded_file.name)
                except:
                    pass  # Fallback gracefully if ingestion fails
                
                st.success(f"✅ Added {uploaded_file.name}")
                st.rerun()
            else:
                st.error(f"Failed to extract text from {uploaded_file.name}")
    else:
        st.info(f"📄 {uploaded_file.name} is already uploaded")
 
# Display welcome message or chat history
if not st.session_state.messages:
    if not st.session_state.uploaded_documents:
        st.markdown('<div class="subtitle">Drop a PDF or PPTX file to get started, or ask me anything</div>', unsafe_allow_html=True)
    else:
        st.markdown(f'<div class="subtitle">{len(st.session_state.uploaded_documents)} document(s) loaded! Ask me anything about them.</div>', unsafe_allow_html=True)
else:
    # Display chat history
    chat_container = st.container()
    with chat_container:
        for message in st.session_state.messages:
            if message["role"] == "user":
                st.markdown(f"""
                    <div class="chat-message user">
                        <div class="message-header">You</div>
                        <div class="message-content">{message["content"]}</div>
                    </div>
                    """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                    <div class="chat-message bot">
                        <div class="message-header">Assistant</div>
                        <div class="message-content">{message["content"]}</div>
                    </div>
                    """, unsafe_allow_html=True)
 
# Add spacing before input
st.markdown("<br>" * 3, unsafe_allow_html=True)
 
# Chat input at the bottom
user_input = st.chat_input("Ask anything")
 
if user_input:
    # Add user message to chat
    st.session_state.messages.append({
        "role": "user",
        "content": user_input,
        "timestamp": datetime.now()
    })
   
    # Get bot response
    with st.spinner("🤔 Thinking..."):
        bot_response = get_bot_response(user_input)
   
    # Add bot response to chat
    st.session_state.messages.append({
        "role": "assistant",
        "content": bot_response,
        "timestamp": datetime.now()
    })
   
    # Rerun to update chat display
    st.rerun()
 
 